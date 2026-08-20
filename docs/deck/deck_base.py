# -*- coding: utf-8 -*-
"""
Bordo 발표덱 생성기 — 공통 디자인 시스템 · 도형 헬퍼.

화이트 베이스 + 파스텔(블루/민트/라벤더/로즈/슬레이트). 주황 계열은 쓰지 않습니다.
모든 도형은 PowerPoint 네이티브 도형이라 발표 직전에 손으로 고칠 수 있습니다.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ─────────────────────────────────────────────── 캔버스
W, H = 13.333, 7.5          # 16:9
ML, MR = 0.86, 0.86         # 좌우 여백
CW = W - ML - MR            # 콘텐츠 폭

FONT = "Noto Sans KR"       # 설치 확인됨
FONT_MONO = "Consolas"

A_NS = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'


def C(h):
    return RGBColor.from_string(h)


# ─────────────────────────────────────────────── 팔레트
INK = C("1B2233")
INK2 = C("47536B")
MUTE = C("8B97AC")
FAINT = C("B8C2D2")
LINE = C("E3E9F1")
WHITE = C("FFFFFF")
SURF = C("F6F8FC")
SURF2 = C("EEF2F8")

BLUE, BLUE_L, BLUE_XL = C("4A7FD4"), C("C9DCF7"), C("EAF2FD")
MINT, MINT_L, MINT_XL = C("2E9C82"), C("C4E9DC"), C("E6F6F1")
LAV, LAV_L, LAV_XL = C("6E5BB8"), C("D8D1F0"), C("EFECFB")
ROSE, ROSE_L, ROSE_XL = C("C4557F"), C("F4CFDF"), C("FCEDF3")
SLATE, SLATE_L, SLATE_XL = C("54708F"), C("D2DDEA"), C("EDF2F8")

ACCENTS = [(BLUE, BLUE_L, BLUE_XL), (MINT, MINT_L, MINT_XL),
           (LAV, LAV_L, LAV_XL), (ROSE, ROSE_L, ROSE_XL)]


# ─────────────────────────────────────────────── 저수준 헬퍼
def _alpha(shp, pct):
    """채우기 투명도. PowerPoint 는 alpha 를 1/1000 % 로 셉니다."""
    sf = shp.fill._xPr.find(qn("a:solidFill"))
    if sf is None or len(sf) == 0:
        return
    clr = sf[0]
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))}))


def _shadow(shp, blur=16, dist=3, alpha=7):
    """카드 그림자. python-pptx 에 API 가 없어 XML 로 직접 넣습니다."""
    spPr = shp._element.spPr
    for e in spPr.findall(qn("a:effectLst")):
        spPr.remove(e)
    xml = (
        "<a:effectLst " + A_NS + ">"
        '<a:outerShdw blurRad="' + str(int(blur * 12700)) + '" dist="'
        + str(int(dist * 12700)) + '" dir="5400000" rotWithShape="0">'
        '<a:srgbClr val="1B2233"><a:alpha val="' + str(int(alpha * 1000))
        + '"/></a:srgbClr></a:outerShdw></a:effectLst>')
    spPr.append(parse_xml(xml))


def _noshadow(shp):
    spPr = shp._element.spPr
    for e in spPr.findall(qn("a:effectLst")):
        spPr.remove(e)
    spPr.append(parse_xml("<a:effectLst " + A_NS + "/>"))


def shape(sl, kind, x, y, w, h, fill=None, line=None, lw=1.0,
          radius=None, shadow=False, alpha=None, rot=None):
    s = sl.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        if alpha is not None:
            _alpha(s, alpha)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = min(0.5, radius / max(0.01, min(w, h)))
        except Exception:
            pass
    if rot is not None:
        s.rotation = rot
    if shadow:
        _shadow(s)
    else:
        _noshadow(s)
    s.text_frame.word_wrap = True
    s.text_frame.margin_left = 0
    s.text_frame.margin_right = 0
    s.text_frame.margin_top = 0
    s.text_frame.margin_bottom = 0
    return s


def card(sl, x, y, w, h, fill=WHITE, line=LINE, radius=0.14, shadow=True, lw=1.0):
    return shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
                 fill=fill, line=line, lw=lw, radius=radius, shadow=shadow)


def box(sl, x, y, w, h, **kw):
    return shape(sl, MSO_SHAPE.RECTANGLE, x, y, w, h, **kw)


def oval(sl, x, y, w, h, **kw):
    return shape(sl, MSO_SHAPE.OVAL, x, y, w, h, **kw)


def pill(sl, x, y, w, h, **kw):
    kw.setdefault("radius", h / 2)
    return shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, **kw)


def _dash(ln):
    ln.line._get_or_add_ln().append(
        parse_xml("<a:prstDash " + A_NS + ' val="sysDash"/>'))


def hline(sl, x, y, w, color=LINE, lw=1.0, dash=False):
    ln = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(lw)
    if dash:
        _dash(ln)
    return ln


def vline(sl, x, y, h, color=LINE, lw=1.0, dash=False):
    ln = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Inches(x), Inches(y), Inches(x), Inches(y + h))
    ln.line.color.rgb = color
    ln.line.width = Pt(lw)
    if dash:
        _dash(ln)
    return ln


def arrow(sl, x, y, w, color=SLATE_L, h=0.17):
    return shape(sl, MSO_SHAPE.RIGHT_ARROW, x, y - h / 2, w, h, fill=color)


def tri(sl, x, y, w, h, fill, rot=0):
    return shape(sl, MSO_SHAPE.ISOSCELES_TRIANGLE, x, y, w, h, fill=fill, rot=rot)


# ─────────────────────────────────────────────── 텍스트
_AL = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
_AN = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def _style(run, size, color, bold, font, spc):
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    rPr = run.font._rPr
    rPr.set("spc", str(int(spc * 5)))
    for tag in ("a:latin", "a:ea", "a:cs"):
        for e in rPr.findall(qn(tag)):
            rPr.remove(e)
    for tag in ("a:latin", "a:ea", "a:cs"):
        rPr.append(parse_xml(
            "<" + tag + " " + A_NS + ' typeface="' + font + '"/>'))


def T(sl, x, y, w, h, text, size=13, color=INK2, bold=False, align="l",
      anchor="t", spc=0, lh=1.35, font=None, space_after=0, target=None):
    font = font or FONT
    if target is None:
        tbx = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        tbx = target
    tf = tbx.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    if target is None:
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
    tf.vertical_anchor = _AN[anchor]
    for i, one in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = _AL[align]
        p.line_spacing = lh
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = one
        _style(r, size, color, bold, font, spc)
    return tbx


def TR(sl, x, y, w, h, runs, align="l", anchor="t", lh=1.35):
    """한 줄 안에서 굵기·색이 갈리는 텍스트. runs=[(글자,size,color,bold[,spc[,font]])]"""
    tbx = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tbx.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = _AN[anchor]
    p = tf.paragraphs[0]
    p.alignment = _AL[align]
    p.line_spacing = lh
    for item in runs:
        txt, size, color, bold = item[0], item[1], item[2], item[3]
        spc = item[4] if len(item) > 4 else 0
        fnt = item[5] if len(item) > 5 else FONT
        r = p.add_run()
        r.text = txt
        _style(r, size, color, bold, fnt, spc)
    return tbx


def fit(sl, s, text, size=12, color=INK, bold=False, align="c", anchor="m",
        spc=0, lh=1.22, font=None):
    """도형 안에 텍스트."""
    return T(sl, 0, 0, 0, 0, text, size, color, bold, align, anchor, spc, lh,
             font, target=s)


# ─────────────────────────────────────────────── 픽토그램
def ic_clock(sl, cx, cy, d, color, lw=1.6):
    oval(sl, cx - d / 2, cy - d / 2, d, d, fill=None, line=color, lw=lw)
    vline(sl, cx, cy - d * 0.27, d * 0.27, color=color, lw=lw)
    hline(sl, cx, cy, d * 0.21, color=color, lw=lw)


def ic_person(sl, cx, cy, d, color):
    oval(sl, cx - d * 0.185, cy - d * 0.42, d * 0.37, d * 0.37, fill=color)
    s = shape(sl, MSO_SHAPE.ROUND_2_SAME_RECTANGLE, cx - d * 0.32,
              cy + d * 0.04, d * 0.64, d * 0.34, fill=color)
    try:
        s.adjustments[0] = 0.5
        s.adjustments[1] = 0.0
    except Exception:
        pass


def ic_doc(sl, cx, cy, d, color, lw=1.5):
    box(sl, cx - d * 0.27, cy - d * 0.37, d * 0.54, d * 0.74, fill=None,
        line=color, lw=lw)
    for i, wf in enumerate((0.32, 0.32, 0.20)):
        hline(sl, cx - d * 0.16, cy - d * 0.15 + i * d * 0.18, d * wf,
              color=color, lw=lw * 0.85)


def ic_bubble(sl, cx, cy, d, color, fill=None, lw=1.5):
    s = shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE, cx - d * 0.35, cy - d * 0.33,
              d * 0.70, d * 0.52, fill=fill, line=color, lw=lw, radius=d * 0.13)
    tri(sl, cx - d * 0.20, cy + d * 0.16, d * 0.15, d * 0.15,
        fill if fill else color, rot=180)
    return s


def ic_badge(sl, cx, cy, d, glyph, fg, bg):
    s = oval(sl, cx - d / 2, cy - d / 2, d, d, fill=bg)
    fit(sl, s, glyph, size=d * 44, color=fg, bold=True)
    return s


def toggle(sl, x, y, w, on, on_color=MINT, off_color=C("CDD6E2")):
    """ON/OFF 스위치 픽토그램. 반환값은 높이."""
    h = w * 0.50
    pill(sl, x, y, w, h, fill=(on_color if on else off_color))
    knob = h * 0.74
    kx = x + (w - knob - h * 0.13) if on else x + h * 0.13
    oval(sl, kx, y + (h - knob) / 2, knob, knob, fill=WHITE)
    return h


# ─────────────────────────────────────────────── 슬라이드 뼈대
def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    return prs


def blank(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    box(sl, 0, 0, W, H, fill=WHITE)
    return sl


PAGE = {"n": 0}


def frame(prs, kicker, title, sub=None, accent=BLUE):
    """공통 헤더·푸터. 본문 시작 y 를 함께 돌려줍니다."""
    sl = blank(prs)
    PAGE["n"] += 1
    box(sl, 0, 0, 0.085, H, fill=accent)
    T(sl, ML, 0.58, CW, 0.24, kicker, 10, accent, True, "l", "t", spc=18)
    T(sl, ML, 0.90, CW * 0.88, 0.60, title, 25, INK, True, "l", "t", lh=1.15)
    y = 1.62
    if sub:
        T(sl, ML, 1.56, CW * 0.84, 0.28, sub, 11.5, MUTE, False, "l", "t")
        y = 1.94
    hline(sl, ML, y, CW, LINE)
    T(sl, ML, H - 0.50, 2.0, 0.22, "Bordo", 9, FAINT, True, "l", "t", spc=30)
    T(sl, W - MR - 1.2, H - 0.50, 1.2, 0.22, "%02d" % PAGE["n"], 9, FAINT,
      True, "r", "t", spc=20)
    return sl, y + 0.32


def source(sl, text, y=None):
    """하단 출처 한 줄."""
    T(sl, ML, y if y is not None else H - 0.92, CW, 0.26, text, 9.5, FAINT,
      False, "l", "t")


def placeholder(sl, x, y, w, h, caption, hint="", accent=BLUE):
    """직접 캡처를 넣을 자리. 점선 박스 + 캡션."""
    s = shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=SURF,
              line=FAINT, lw=1.0, radius=0.12)
    s.line._get_or_add_ln().append(
        parse_xml("<a:prstDash " + A_NS + ' val="dash"/>'))
    ic_badge(sl, x + w / 2, y + h / 2 - 0.30, 0.44, "+", accent, WHITE)
    T(sl, x + 0.16, y + h / 2 - 0.02, w - 0.32, 0.26, caption, 11, INK2,
      True, "c", "t")
    if hint:
        T(sl, x + 0.16, y + h / 2 + 0.26, w - 0.32, 0.5, hint, 9.5, MUTE,
          False, "c", "t", lh=1.3)
    return s
